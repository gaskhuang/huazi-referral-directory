#!/usr/bin/env python3
"""從公開匯出網址抓簡報，攤平成 data/slides.json。完全不需要任何憑證。

簡報的共用設定是「知道連結的人都能存取」，所以 Google 的匯出端點
    https://docs.google.com/presentation/d/<ID>/export/pptx
不帶任何 token 就能下載。抓下來後用 python-pptx 讀版面座標與文字，
輸出格式跟 scripts/extract.py 一致，後面的 parse.py / build.py 不用改。
"""
import json, os, sys, urllib.request

DECK_ID = "1j9lY9x0Iswm-aR6bdR7NBVLIFX7g100DEz7DNt4hXnw"
EXPORT = f"https://docs.google.com/presentation/d/{DECK_ID}/export/pptx"
PPTX = "data/raw/deck.pptx"


def download():
    os.makedirs("data/raw", exist_ok=True)
    print(f"→ 下載公開匯出 {DECK_ID}")
    req = urllib.request.Request(EXPORT, headers={"User-Agent": "Mozilla/5.0"})
    with urllib.request.urlopen(req, timeout=300) as resp, open(PPTX, "wb") as f:
        total = 0
        while True:
            chunk = resp.read(1 << 20)
            if not chunk:
                break
            f.write(chunk)
            total += len(chunk)
    if total < 100_000:
        sys.exit(f"✗ 下載的檔案只有 {total} bytes，簡報可能已改成不公開")
    print(f"→ 下載完成 {total / 1024 / 1024:.1f} MB")


def group_xfrm(sh):
    """群組的 a:off / a:ext / a:chOff / a:chExt（python-pptx 沒有直接的屬性）"""
    from pptx.oxml.ns import qn
    gp = sh._element.find(qn("p:grpSpPr"))
    xf = gp.find(qn("a:xfrm")) if gp is not None else None
    if xf is None:
        return (sh.left or 0, sh.top or 0, sh.width or 0, sh.height or 0, 0, 0, 0, 0)

    def pair(tag, a, b):
        n = xf.find(qn(tag))
        return (int(n.get(a)), int(n.get(b))) if n is not None else (0, 0)

    ox, oy = pair("a:off", "x", "y")
    ex, ey = pair("a:ext", "cx", "cy")
    cox, coy = pair("a:chOff", "x", "y")
    cex, cey = pair("a:chExt", "cx", "cy")
    return ox, oy, ex, ey, cox, coy, cex, cey


def walk(shapes, out, dx=0, dy=0, sx=1.0, sy=1.0):
    """遞迴展開群組，把子圖形的座標換算回投影片的絕對座標"""
    for sh in shapes:
        if sh.shape_type == 6 and hasattr(sh, "shapes"):        # GROUP
            gx, gy, gw, gh, cx, cy, cw, ch = group_xfrm(sh)
            nsx = (gw / cw) if cw else 1.0
            nsy = (gh / ch) if ch else 1.0
            walk(sh.shapes, out,
                 dx + (gx - cx * nsx) * sx, dy + (gy - cy * nsy) * sy,
                 sx * nsx, sy * nsy)
            continue

        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t:
                out.append({
                    "y": int(dy + (sh.top or 0) * sy),
                    "x": int(dx + (sh.left or 0) * sx),
                    "w": int((sh.width or 0) * sx),
                    "h": int((sh.height or 0) * sy),
                    "t": t,
                })
        elif getattr(sh, "has_table", False) and sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        out.append({
                            "y": int(dy + (sh.top or 0) * sy),
                            "x": int(dx + (sh.left or 0) * sx),
                            "w": int((sh.width or 0) * sx),
                            "h": int((sh.height or 0) * sy),
                            "t": t,
                        })


def main():
    if "--skip-download" not in sys.argv or not os.path.exists(PPTX):
        download()

    try:
        from pptx import Presentation
    except ImportError:
        sys.exit("✗ 需要 python-pptx：pip install python-pptx")

    prs = Presentation(PPTX)
    slides = []
    for i, s in enumerate(list(prs.slides)):
        boxes = []
        walk(s.shapes, boxes)
        pics = sum(1 for sh in s.shapes if not sh.has_text_frame)
        boxes.sort(key=lambda b: (b["y"], b["x"]))
        slides.append({"i": i, "pics": pics, "boxes": boxes})

    os.makedirs("data", exist_ok=True)
    json.dump(slides, open("data/slides.json", "w", encoding="utf-8"),
              ensure_ascii=False, indent=1)
    print(f"→ 攤平 {len(slides)} 頁，"
          f"文字框合計 {sum(len(s['boxes']) for s in slides)}，"
          f"無文字頁 {sum(1 for s in slides if not s['boxes'])}")


if __name__ == "__main__":
    main()
